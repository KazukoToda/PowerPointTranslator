from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os

def extract_text_from_pptx(pptx_path):
    """
    Extract text and formatting information from a PowerPoint file
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        
    Returns:
        list: List of slides, each containing a list of (text, shape_info) tuples
    """
    prs = Presentation(pptx_path)
    slides_data = []
    
    for slide in prs.slides:
        slide_elements = []
        
        for shape in slide.shapes:
            # Process text-containing shapes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    
                    # Store text and shape information
                    shape_info = {
                        'type': 'text',
                        'shape_id': shape.shape_id,
                        'x': shape.left,
                        'y': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'font_size': _get_font_size(paragraph),
                        'font_name': _get_font_name(paragraph),
                        'color': _get_text_color(paragraph),
                        'bold': _is_bold(paragraph),
                        'italic': _is_italic(paragraph)
                    }
                    
                    slide_elements.append((text, shape_info))
            
            # Store information about images
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info = {
                    'type': 'image',
                    'shape_id': shape.shape_id,
                    'x': shape.left,
                    'y': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                }
                slide_elements.append(("", shape_info))
                
            # Store information about tables
            elif shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for paragraph in cell.text_frame.paragraphs:
                            text = paragraph.text
                            
                            shape_info = {
                                'type': 'table_cell',
                                'shape_id': shape.shape_id,
                                'row': row_idx,
                                'column': col_idx,
                                'font_size': _get_font_size(paragraph),
                                'font_name': _get_font_name(paragraph),
                                'color': _get_text_color(paragraph)
                            }
                            
                            slide_elements.append((text, shape_info))
        
        slides_data.append(slide_elements)
    
    return slides_data

def create_translated_pptx(input_path, output_path, translated_slides):
    """
    Create a new PowerPoint file with translated text
    
    Args:
        input_path (str): Path to the original PowerPoint file
        output_path (str): Path to save the translated PowerPoint file
        translated_slides (list): List of slides with translated text
    """
    # Load the original presentation
    prs = Presentation(input_path)
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= len(translated_slides):
            continue
            
        translated_elements = translated_slides[slide_idx]
        shape_texts = {}
        
        # Group translated texts by shape_id
        for text, shape_info in translated_elements:
            shape_id = shape_info.get('shape_id')
            
            if shape_id is not None:
                if shape_id not in shape_texts:
                    shape_texts[shape_id] = []
                
                if shape_info.get('type') == 'table_cell':
                    # For table cells, store with row and column
                    row = shape_info.get('row', 0)
                    col = shape_info.get('column', 0)
                    shape_texts[shape_id].append((text, 'table_cell', row, col))
                else:
                    # For text shapes, append to list of paragraphs
                    shape_texts[shape_id].append((text, 'text'))
        
        # Apply translated text to shapes
        for shape in slide.shapes:
            if shape.shape_id in shape_texts:
                elements = shape_texts[shape.shape_id]
                
                if shape.has_text_frame:
                    # Clear existing text
                    current_paragraphs = len(shape.text_frame.paragraphs)
                    text_elements = [e for e in elements if e[1] == 'text']
                    
                    # Apply translated text
                    for i, (text, _) in enumerate(text_elements):
                        if i < current_paragraphs:
                            p = shape.text_frame.paragraphs[i]
                            if p.runs:
                                p.runs[0].text = text
                            else:
                                p.text = text
                        else:
                            p = shape.text_frame.add_paragraph()
                            p.text = text
                
                elif shape.has_table:
                    table_elements = [e for e in elements if e[1] == 'table_cell']
                    
                    # Apply translated text to table cells
                    for text, _, row, col in table_elements:
                        if row < len(shape.table.rows) and col < len(shape.table.rows[row].cells):
                            cell = shape.table.rows[row].cells[col]
                            if cell.text_frame.paragraphs:
                                cell.text_frame.paragraphs[0].text = text
    
    # Save the presentation
    prs.save(output_path)
    return output_path

def _get_font_size(paragraph):
    """Helper function to get font size from paragraph"""
    if paragraph.runs and paragraph.runs[0].font.size:
        return paragraph.runs[0].font.size
    return None

def _get_font_name(paragraph):
    """Helper function to get font name from paragraph"""
    if paragraph.runs and paragraph.runs[0].font.name:
        return paragraph.runs[0].font.name
    return None

def _get_text_color(paragraph):
    """Helper function to get text color from paragraph"""
    if paragraph.runs and paragraph.runs[0].font.color.rgb:
        rgb = paragraph.runs[0].font.color.rgb
        return f"{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
    return None

def _is_bold(paragraph):
    """Helper function to check if text is bold"""
    if paragraph.runs and paragraph.runs[0].font.bold:
        return paragraph.runs[0].font.bold
    return False

def _is_italic(paragraph):
    """Helper function to check if text is italic"""
    if paragraph.runs and paragraph.runs[0].font.italic:
        return paragraph.runs[0].font.italic
    return False